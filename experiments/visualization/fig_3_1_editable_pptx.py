"""Build an editable PowerPoint assembly of Fig. 3-1.

The scientific image panels remain lossless cropped raster evidence. Titles,
captions, arrows, separators and labels are native PowerPoint objects.
"""

from __future__ import annotations

import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
FIG_DIR = ROOT / "experiments" / "visualization" / "outputs" / "FIG-3-1"
SOURCE = FIG_DIR / "fig_3_1.png"
ASSET_DIR = FIG_DIR / "editable_pptx_assets"
OUTPUT = FIG_DIR / "fig_3_1_editable.pptx"
METADATA = FIG_DIR / "metadata.json"

SLIDE_W = 13.333
SLIDE_H = 7.5
INK = RGBColor(38, 50, 56)
MUTED = RGBColor(102, 117, 124)
LIGHT = RGBColor(217, 222, 226)
BLUE = RGBColor(47, 111, 163)
AMBER = RGBColor(216, 138, 35)
TEAL = RGBColor(42, 157, 143)
GREEN = RGBColor(75, 154, 118)
WHITE = RGBColor(255, 255, 255)


# Pixel regions exclude most figure annotations so the editable slide supplies
# consistent native text and arrows. Their positions are tied to the 600 dpi
# master PNG produced by fig_3_1_overall_workflow.py.
CROPS = {
    "dom": (149, 382, 660, 1038),
    "input_cloud": (725, 382, 1236, 1038),
    "tiles": (149, 1250, 1236, 2021),
    "candidate_a": (1494, 556, 2074, 1211),
    "candidate_b": (2263, 556, 2843, 1211),
    "fused": (1724, 1529, 2614, 2011),
    "mapping": (3155, 440, 4134, 941),
    "integration": (3155, 1211, 4134, 1867),
}


def crop_assets() -> dict[str, Path]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    with Image.open(SOURCE) as image:
        result = {}
        for name, bounds in CROPS.items():
            target = ASSET_DIR / f"{name}.png"
            image.crop(bounds).save(target, dpi=(600, 600))
            result[name] = target
    return result


def add_text(slide, text: str, x: float, y: float, w: float, h: float, *, size: float, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    return box


def add_rule(slide, x: float, y: float, w: float) -> None:
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.012))
    rule.fill.solid()
    rule.fill.fore_color.rgb = LIGHT
    rule.line.fill.background()


def add_arrow(slide, x: float, y: float, w: float, h: float, *, vertical=False) -> None:
    shape_type = MSO_SHAPE.DOWN_ARROW if vertical else MSO_SHAPE.RIGHT_ARROW
    arrow = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = MUTED
    arrow.line.fill.background()


def add_label(slide, text: str, x: float, y: float, color) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(0.66), Inches(0.24))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.fill.transparency = 9
    shape.line.fill.background()
    add_text(slide, text, x + 0.03, y + 0.01, 0.60, 0.20, size=7.1, color=color, bold=True, align=PP_ALIGN.CENTER)


def add_picture(slide, image: Path, x: float, y: float, w: float, h: float) -> None:
    slide.shapes.add_picture(str(image), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def add_band_heading(slide, title: str, subtitle: str, x: float, w: float) -> None:
    add_text(slide, title, x, 0.25, w, 0.30, size=12.0, bold=True)
    add_text(slide, subtitle, x, 0.58, w, 0.20, size=7.2, color=MUTED)
    add_rule(slide, x, 0.84, w)


def main() -> None:
    if not SOURCE.exists():
        raise FileNotFoundError(f"Missing Fig. 3-1 master PNG: {SOURCE}")
    assets = crop_assets()
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_W)
    presentation.slide_height = Inches(SLIDE_H)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # Three editable workflow bands.
    add_band_heading(slide, "DOM Preparation & Adaptive Tiling", "Co-registered DOM + original local point cloud", 0.35, 3.42)
    add_band_heading(slide, "Tile-wise Segmentation & Fusion", "Per-tile inference -> correlation clustering", 4.66, 4.02)
    add_band_heading(slide, "Point-Cloud Mapping & 2.5D Measurement", "Fused footprint + local point cloud -> 2.5D measurement", 9.57, 3.42)

    add_text(slide, "DOM local view", 0.48, 1.03, 1.38, 0.22, size=8.0, align=PP_ALIGN.CENTER)
    add_text(slide, "Original point cloud", 2.15, 1.03, 1.38, 0.22, size=8.0, align=PP_ALIGN.CENTER)
    add_picture(slide, assets["dom"], 0.48, 1.27, 1.38, 1.76)
    add_picture(slide, assets["input_cloud"], 2.15, 1.27, 1.38, 1.76)
    add_arrow(slide, 2.02, 3.18, 0.10, 0.32, vertical=True)
    add_text(slide, "Adaptive overlapping tiles", 0.58, 3.56, 2.88, 0.22, size=8.0, align=PP_ALIGN.CENTER)
    add_picture(slide, assets["tiles"], 0.48, 3.80, 3.05, 2.15)
    add_label(slide, "Tile B", 0.60, 3.95, BLUE)
    add_label(slide, "Tile A", 0.60, 5.65, AMBER)

    yolo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.53), Inches(0.98), Inches(1.90), Inches(0.30))
    yolo.fill.solid()
    yolo.fill.fore_color.rgb = WHITE
    yolo.line.color.rgb = INK
    yolo.line.width = Pt(1.0)
    add_text(slide, "YOLO11m-seg", 5.62, 0.97, 1.72, 0.16, size=8.0, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "tile-wise instance segmentation", 5.62, 1.12, 1.72, 0.13, size=5.6, color=MUTED, align=PP_ALIGN.CENTER)

    add_text(slide, "Tile A candidate", 4.70, 1.42, 1.65, 0.22, size=8.0, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Tile B candidate", 6.96, 1.42, 1.65, 0.22, size=8.0, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_picture(slide, assets["candidate_a"], 4.70, 1.66, 1.65, 1.86)
    add_picture(slide, assets["candidate_b"], 6.96, 1.66, 1.65, 1.86)
    add_arrow(slide, 5.73, 3.66, 0.10, 0.26, vertical=True)
    add_arrow(slide, 7.72, 3.66, 0.10, 0.26, vertical=True)
    add_text(slide, "correlation clustering", 5.36, 3.94, 2.58, 0.16, size=7.4, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "world-coordinate mask IoU + centroid distance", 5.36, 4.12, 2.58, 0.16, size=6.2, color=MUTED, align=PP_ALIGN.CENTER)
    add_arrow(slide, 6.59, 4.34, 0.10, 0.22, vertical=True)
    add_picture(slide, assets["fused"], 5.53, 4.64, 2.26, 1.22)
    add_text(slide, "Fused rock footprint", 5.53, 5.90, 2.26, 0.20, size=8.0, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, "Local point-cloud mapping (top view)", 9.67, 1.03, 3.22, 0.22, size=8.0, bold=True, align=PP_ALIGN.CENTER)
    add_picture(slide, assets["mapping"], 9.78, 1.27, 3.00, 1.54)
    add_arrow(slide, 11.23, 2.95, 0.10, 0.28, vertical=True)
    add_text(slide, "2.5D integration cells (top view)", 9.67, 3.40, 3.22, 0.22, size=8.0, bold=True, align=PP_ALIGN.CENTER)
    add_picture(slide, assets["integration"], 9.78, 3.66, 3.00, 2.01)

    add_arrow(slide, 3.90, 3.25, 0.35, 0.16)
    add_arrow(slide, 8.90, 3.25, 0.35, 0.16)
    presentation.save(OUTPUT)

    metadata = json.loads(METADATA.read_text(encoding="utf-8"))
    metadata["editable_presentation"] = {
        "path": str(OUTPUT.relative_to(ROOT)),
        "script": str(Path(__file__).relative_to(ROOT)),
        "format": "PowerPoint .pptx",
        "editable_objects": "Text, arrows, labels, rules and image-panel placement are native slide objects; evidence panels are embedded high-resolution raster images.",
    }
    outputs = metadata.setdefault("outputs", [])
    for name in (OUTPUT.name, Path(__file__).name):
        if name not in outputs:
            outputs.append(name)
    METADATA.write_text(json.dumps(metadata, ensure_ascii=False, indent=2), encoding="utf-8")
    print(OUTPUT)


if __name__ == "__main__":
    main()
